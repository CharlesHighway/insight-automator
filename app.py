import streamlit as st
import pandas as pd
from pptx import Presentation
from pptx.util import Inches
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
import os

st.title("📊 Insight Automator (Demo)")

# Toggle between mock / GPT
USE_MOCK = st.checkbox("Use mock mode (no GPT calls)", value=True)

api_key = None
if not USE_MOCK:
    api_key = st.text_input("Enter your OpenAI API Key", type="password")
    if api_key:
        from openai import OpenAI
        client = OpenAI(api_key=api_key)

uploaded_file = st.file_uploader("Upload survey data (CSV)", type=["csv"])
template_file = st.file_uploader("Upload PowerPoint template", type=["pptx"])

if uploaded_file and template_file:
    df = pd.read_csv(uploaded_file)
    st.write("### Data Preview", df.head())

    # Crosstab pairs (hard-coded for demo; we can make this user-selectable later)
    CROSSTAB_PAIRS = [("Gender", "Preference"), ("Region", "Preference")]

    def run_crosstab(df, row_var, col_var):
        counts = pd.crosstab(df[row_var], df[col_var])
        percents = pd.crosstab(df[row_var], df[col_var], normalize="columns") * 100
        return {"counts": counts, "percents": percents.round(1)}

    def generate_insights(var_name, crosstab):
        if USE_MOCK:
            return f"Insights for {var_name}", [
                f"{var_name} – point 1 (mock)",
                f"{var_name} – point 2 (mock)",
                f"{var_name} – point 3 (mock)"
            ]
        else:
            if not api_key:
                return f"{var_name} (Error)", ["⚠️ Please enter API key"]
            prompt = f"""
            You are a market research consultant.
            Crosstab analysis: {var_name}
            {crosstab.to_dict()}
            Write 1 headline + 3 key points.
            """
            response = client.chat.completions.create(
                model="gpt-4o-mini",
                messages=[
                    {"role":"system","content":"You are a senior insights consultant."},
                    {"role":"user","content":prompt}
                ]
            )
            text = response.choices[0].message.content
            lines = text.split("\\n")
            return lines[0], lines[1:4]

    # Build PowerPoint
    prs = Presentation(template_file)

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Market Research Report"

    for row_var, col_var in CROSSTAB_PAIRS:
        ct = run_crosstab(df, row_var, col_var)
        title, bullets = generate_insights(f"{row_var} vs {col_var}", ct["percents"])

        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        tf = slide.placeholders[1].text_frame
        tf.clear()
        for b in bullets:
            tf.add_paragraph().text = b

        data = CategoryChartData()
        data.categories = list(ct["percents"].index)
        for col in ct["percents"].columns:
            data.add_series(col, list(ct["percents"][col]))
        slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
                               Inches(0.5), Inches(2.5), Inches(8), Inches(3.5), data)

    output_file = "insight_report.pptx"
    prs.save(output_file)
    st.success("✅ Report generated!")

    with open(output_file, "rb") as f:
        st.download_button("⬇️ Download PowerPoint", f, file_name=output_file)
